# -*- coding: utf-8 -*-
"""扫描日期文件夹中的 PDF 文献，生成 PPT 与 HTML 检索页。"""

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from paper_zone import classify_cloud_edge_research, enrich_paper_item

SITE_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = SITE_DIR
ROOT = SITE_DIR.parent.parent  # D:\刚需 — PDF date folders
DATE_FOLDERS = ["4.9", "4.17", "4.24", "5.1", "5.8", "5.16", "5.23", "5.31", "6.6", "6.13", "6.20", "6.27", "7.10", "7.17", "7.24", "7.31", "8.8", "8.16", "8.23"]

# 一句话中文概括（按文件名 stem 匹配）
SUMMARIES = {
    "Beamforming_and_Device_Selection_Design_in_Federated_Learning_With_Over-the-Air_Aggregation": "研究联邦学习空中聚合场景下，如何联合优化波束成形与参与设备选择以提升训练效率。",
    "Cloud-Edge_Collaborative_Depression_Detection_Using_Negative_Emotion_Recognition_and_Cross-Scale_Facial_Feature_Analysis": "提出云边协作抑郁症检测方案，结合负面情绪识别与跨尺度面部特征分析实现高效筛查。",
    "Massive_Digital_Over-the-Air_Computation_for_Communication-Efficient_Federated_Edge_Learning": "利用大规模数字空中计算降低联邦边缘学习的通信开销，加速模型聚合。",
    "Optimal_Client_Sampling_in_Federated_Learning_with_Client-level_Heterogeneous_Differential_Privacy": "在客户端异构差分隐私约束下，设计最优客户端采样策略以平衡隐私与模型精度。",
    "Primal-Dual-Based_Computation_Offloading_Method_for_Energy-Aware_Cloud-Edge_Collaboration": "基于原始-对偶优化设计能量感知的云边协作计算卸载方法，降低系统能耗。",
    "AI-Enhanced_Cloud-Edge-Terminal_Collaborative_Network_Survey_Applications_and_Future_Directions": "综述 AI 增强的云-边-端协同网络架构、典型应用与未来研究方向。",
    "Reinforcement_Learning_for_Task_Placement_in_Collaborative_Cloud-_Edge_Computing": "用强化学习解决云边协同计算中的任务放置问题，动态适配资源与负载变化。",
    "Transcoding-Enabled_CloudEdgeTerminal_Collaborative_Video_Caching_in_Heterogeneous_IoT_Networks_An_Online_Learning_Approach_With_Time-Varying_Information": "面向异构 IoT 网络，提出支持转码的云边端协同视频缓存在线学习方案。",
    "Delay-Aware_Cooperative_Task_Offloading_for_Multi-UAV_Enabled_Edge-Cloud_Computing": "在多无人机边缘-云系统中，考虑时延约束的协作任务卸载与调度优化。",
    "Efficient_Online_Computing_Offloading_for_Budget-_Constrained_Cloud-Edge_Collaborative_Video_Streaming_Systems": "在预算约束下，为云边协同视频流系统设计高效的在线计算卸载策略。",
    "Time-Slotted_Task_Offloading_and_Resource_Allocation_for_Cloud-Edge-End_Cooperative_Computing_Networks": "针对云-边-端协同网络，研究时隙化任务卸载与资源联合分配机制。",
    "Online_Collaborative_Resource_Allocation_and_Task_Offloading_for_Multi-Access_Edge_Computing": "面向多接入边缘计算，提出在线协作式资源分配与任务卸载框架。",
    "Coding_Assisted_Cloud-Edge_Collaborative_Computing": "利用 MDS 编码辅助云边协同计算，缓解慢节点（Straggler）拖慢整体进度的问题。",
    "EdgeShard_Efficient_LLM_Inference_via_Collaborative_Edge_Computing": "通过协作边缘计算实现大语言模型（LLM）的高效分片推理，降低单节点负载。",
    "Collaborative_Service_Caching_Task_Offloading_and_Resource_Allocation_in_Caching-Assisted_Mobile_Edge_Computing": "在缓存辅助移动边缘计算中，联合优化服务缓存、任务卸载与资源分配。",
    "Truthful_Online_Double_Auction-Based_Resource_Allocation_Mechanisms_for_Partial_Computation_Offloading_in_Collaborative_Edge_Computing": "设计真实在线双拍卖机制，为协作边缘计算中的部分计算卸载分配资源。",
    "Computation_Offloading_Optimization_for_UAV-Based_Cloud-Edge_Collaborative_Task_Scheduling_Strategy": "优化无人机云边协同场景下的计算卸载与任务调度策略。",
    "Flocky_Decentralized_Intent-Based_Edge_Orchestration_Using_Open_Application_Model": "提出去中心化意图驱动边缘编排框架 Flocky，基于 OAM 实现大规模边缘服务部署。",
    "Graph_Neural_Networks_and_Deep_Reinforcement_Learning-Based_Resource_Allocation_for_V2X_Communications": "结合图神经网络与深度强化学习，优化车联网（V2X）通信中的资源分配。",
    "KCES_A_Workflow_Containerization_Scheduling_Scheme_Under_Cloud-Edge_Collaboration_Framework": "在云边协作框架下提出 KCES 工作流容器化调度方案，提升任务编排效率。",
    "Blockchain-Secured_Online_Edge_Collaboration_in_IoT_Integrating_Convex_Optimization_and_Learning_Approach": "融合区块链、凸优化与学习算法，保障 IoT 在线边缘协作的安全与效率。",
    "SDG-CDA_Stackelberg_Differential_Games_and_Combinatorial_Double_Auctions-Based_Pricing_Mechanism_in_CloudEdge_Environment": "在云边环境中，基于 Stackelberg 微分博弈与组合双拍卖设计定价与资源交易机制。",
    "Cloud-EdgeEnd_Collaborative_Task_Offloading_in_Vehicular_Edge_Networks_A_Multilayer_Deep_Reinforcement_Learning_Approach": "面向车载边缘网络，用多层深度强化学习实现云-边-端协同任务卸载。",
    "Task_Offloading_and_Resource_Optimization_Based_on_Dependency-Aware_Graph_and_Collaborative_Deep_Reinforcement_Learning_in_Mobile_Edge_Computing": "基于依赖感知任务图与协作深度强化学习，优化移动边缘计算中的卸载与资源分配。",
    "Cacomp_A_Cloud-Assisted_Collaborative_Deep_Learning_Compiler_Framework_for_DNN_Tasks_on_Edge": "提出云辅助协作深度学习编译框架 Cacomp，优化边端 DNN 任务的部署与执行。",
    "Game-Theoretic_Bandwidth_Allocation_and_Task_Offloading_in_CloudEdge_Collaboration": "用博弈论方法联合优化云边协作中的带宽分配与任务卸载决策。",
    "Two-Stage_Generative_Color_Calibration_for_Drone_Photography_With_Cloud-Edge_Collaboration": "云边协作下无人机摄影的两阶段生成式色彩校准，提升航拍图像质量。",
    "A_Privacy-Preserving_IoT_Data_Access_Control_Scheme_for_Cloud-Edge_Computing": "提出云边计算场景下保护隐私的 IoT 数据访问控制方案，兼顾安全与协作效率。",
    "Differentially_Private_Federated_Tensor_Completion_for_CloudEdge_Collaborative_AIoT_Data_Prediction": "在云边协同 AIoT 数据预测中，采用差分隐私联邦张量补全保护数据隐私。",
    "MEDIA_An_Incremental_DNN_Based_Computation_Offloading_for_Collaborative_Cloud-Edge_Computing": "提出 MEDIA 增量 DNN 计算卸载框架，适配协作云边计算中的模型推理任务。",
    "CloudEdge_Collaborative_SFC_Mapping_for_Industrial_IoT_Using_Deep_Reinforcement_Learning": "面向工业 IoT，用深度强化学习实现云边协同服务功能链（SFC）映射与部署。",
    "Crossport_A_Cloud-Edge-End_Microservice_Architecture_for_Collaborative_Rendering_in_Metaverse_Services": "提出 Crossport 云-边-端微服务架构，支撑元宇宙场景下的协作渲染服务。",
    "Joint_Optimization_of_Dynamic_Service_Selection_and_Request_Routing_in_Cloud-Edge_Collaborative_Environments": "在云边协同环境中，联合优化动态服务选择与请求路由，提升服务交付效率。",
    "A_Differential_Privacy_Based_Task_Offloading_Algorithm_for_Vehicular_Edge_Computing": "面向车载边缘计算，设计基于差分隐私的任务卸载算法，在保护数据隐私的同时优化卸载决策。",
    "A_reliability_and_latency_driven_task_allocation_framework_for_workflow_applications_in_the_edge_hub_cloud_continuum": "针对边-枢纽-云连续体中的工作流应用，提出兼顾可靠性与时延的任务分配框架。",
    "A_Trust_Based_Computation_Offloading_Framework_in_Mobile_Cloud_Edge_Computing_Networks": "在移动云边计算网络中，基于信任机制设计计算卸载框架，提升协作安全性与效率。",
    "Accelerating_AI_Generated_Content_Collaborative_Inference_Via_Transfer_Reinforcement_Learning_in_Dynamic_Edge_Networks": "在动态边缘网络中，利用迁移强化学习加速 AIGC 协同推理，提升生成式 AI 服务响应速度。",
    "AoI_Oriented_Computation_Offloading_and_Resource_Allocation_for_End_Edge_Cloud_Computing_Systems": "面向端-边-云计算系统，以信息年龄（AoI）为导向联合优化计算卸载与资源分配。",
    "Cloud_Edge_Collaborative_Service_Architecture_With_Large_Tiny_Models_Based_on_Deep_Reinforcement_Learning": "提出基于大小模型协同的云边服务架构，用深度强化学习实现模型部署与服务调度。",
    "Cloud_Edge_Collaboration_for_Industrial_Internet_of_Things_Scalable_Neurocomputing_and_Rolling_Horizon_Optimization": "面向工业物联网，结合可扩展神经计算与滚动时域优化实现云边协作智能决策。",
    "An_Optimal_Transport-Based_Federated_Reinforcement_Learning_Approach_for_Resource_Allocation_in_CloudEdge_Collaborative_IoT": "提出 OTFAC 方案，将最优传输与联邦 Actor-Critic 结合，实现云边 IoT 资源分配的时延与能耗联合优化。",
    "Cloud-Edge-End_Collaborative_Computing-Enabled_Intelligent_Sharding_Blockchain_for_Industrial_IoT_Based_on_PPO_Approach": "面向工业 IoT，构建云边端三层协同的智能分片区块链，用 PPO 动态切换共识协议。",
    "Distributed_Resource_Allocation_and_Coordinated_Scheduling_for_End_Edge_Cloud_Collaborative_Computing": "提出端边云分布式资源分配与混合卸载框架，同时支持全卸载与部分卸载，提升系统效用与可扩展性。",
    "Dynamic_Caching_Dependency_Aware_Task_Offloading_in_Mobile_Edge_Computing": "针对 DAG 依赖型移动应用，联合动态缓存与双层优先级调度，用 DRL 优化 MEC 任务卸载。",
    "Edge_Cloud_Continuum_Orchestration_of_Critical_Services_A_Smart_City_Approach": "扩展 Kubernetes 实现边云连续体关键服务编排，自定义实时调度与最快响应负载均衡，时延降低约 80%。",
    "Edge_Collaboration_Enabled_Online_Energy_Optimization_for_Satellite_Assisted_Internet_of_Things_A_Lyapunov_Based_Learnin": "在卫星辅助 IoT 场景中，结合 Lyapunov 优化与 MADRL 实现在线能耗优化与队列稳定。",
    "Efficient_and_Fault_Tolerant_Data_Stream_Processing_With_Uncertain_Data_Rates_in_Serverless_Edge_Computing": "面向无服务器边计算中的 AI 数据流，两阶段算法联合容错部署与在线调度，应对不确定数据速率。",
    "Efficient_End_Edge_Cloud_Task_Offloading_in_6G_Networks_Based_on_Multiagent_DRL": "在 6G 端边云网络中，用 MADDPG 多智能体深度强化学习实现分布式任务卸载与资源优化。",
    "Enhancing_Edge_Cloud_Collaboration_With_Blockchain_Assisted_Digital_Twin_Intelligence_Offloading_Scheme_(GOTTEN)": "提出 GOTTEN 方案，融合区块链与数字孪生网络，实现边云协作智能卸载与资源可验证交易。",
    "When_Crowd_Meets_Big_Video_Data_Cloud-Edge_Collaborative_Transcoding_for_Personal_Livecast": "针对个人直播场景，提出云边协同转码框架，通过观众选择与拍卖激励降低转码成本。",
    # --- 7.17 ---
    "Collaborative_Offloading_for_Interacting_Users_in_Cloud-Edge-Terminal_Networks": "面向云边端交互用户场景，用 SAC 双头混合动作与迁移学习实现协同卸载，兼顾时延与价格成本。",
    "Cost-Optimized_Task_Offloading_for_Dependent_Applications_in_Collaborative_Edge_and_Cloud_Computing": "针对依赖子任务应用，提出 GSPAE（遗传+退火 PSO+自编码器）优化边云协同卸载总成本。",
    "Distributed_Optimization_of_Task_Offloading_and_Resource_Allocation_for_MEC_With_Multifactorial_Uncertainty": "在多因素不确定性下，用 DGOLB 博弈与负载均衡实现 MEC 分布式任务卸载与资源分配。",
    "Digital_Twin_Freshness_Maximization_in_Edge_Computing": "研究边缘计算中数字孪生新鲜度最大化问题，给出近似算法与在线竞争算法。",
    "FCER_A_Federated_Cloud-Edge_Recommendation_Framework_With_Cluster-Based_Edge_Selection": "提出 FCER 联邦云边推荐框架，结合 LDP 预聚类与数据质量驱动的边缘选择。",
    "Towards_Optimal_Customized_Architecture_for_Heterogeneous_Federated_Learning_With_Contrastive_Cloud-Edge_Model_Decoupling": "提出 FedCMD，通过对比学习实现云边模型解耦与异构联邦的动态个性化架构。",
    "Improved_Multiscale_Networks_and_Collaborative_Edge_Computing_for_Precise_and_Real-Time_Tool_Fault_Diagnosis_in_Industrial_IoT": "提出 IMSNet（多尺度卷积+LSTM+多头注意力）与云边端协作卸载 CEDC，实现工业 IoT 刀具故障高精度实时诊断。",
    "CloudEdgeEnd_CFMMIMO_for_UAV_Swarms_Integrated_Sensing_and_Robust_Interference_Detection_in_Perception_Networks": "面向感知网络，研究云边端 CF-mMIMO 无人机群通感一体与鲁棒干扰检测。",
    "Revolutionizing_Turn-by-Turn_Navigation_With_Cloud-Edge_Deep_Learning": "用 Transformer、MoE 与云边协同推理，实现逐向导航语音指令的云边深度学习方案。",
    "Trinity_of_SafetyQualityEfficiency_CloudEdgeDevice_Collaborative_Monitoring_for_Manufacturing_Systems_With_Industrial_Validation": "提出安全—质量—效率三位一体的云边端协同制造监测框架，并完成工业现场验证。",
    # --- 7.24（文件名含空格/en-dash，键须与 pdf.stem 完全一致）---
    "Advancements in RIS-Assisted UAV for Empowering Multiaccess Edge Computing A Survey": "综述 RIS 辅助无人机赋能多接入边缘计算（MEC）的架构、关键技术与开放挑战。",
    "Advancing Traffic Resource Scheduling With Cloud–Edge Collaboration A Virtualized Digital Twin Perspective": "从虚拟化数字孪生视角，研究云边协同下的交通资源调度优化。",
    "An Energy-Efficient Intrusion Detection Offloading Based on DNN for Edge Computing": "面向边缘计算提出基于 DNN 的节能入侵检测卸载方案，降低检测能耗。",
    "An Entropy-Based Privacy-Preserving Federated Deep Reinforcement Learning Framework for Task Offloading in Vehicular Edg": "提出基于熵的隐私保护联邦深度强化学习框架，优化车载边缘任务卸载。",
    "An Improved Proximal Optimization Method for Cloud–Edge–End Collaboration in Railway Surveillance Networks": "改进近端优化方法，支撑铁路监控网络中的云边端协同计算。",
    "An IoT-Oriented Image Retrieval Scheme Based on Multifeature Fusion for Cloud–Edge Environments": "面向 IoT 的云边环境多特征融合图像检索方案，提升检索效率与准确性。",
    "Blocked Job Scheduling and Redundant Computing Resource Allocation in Edge Computing Systems": "研究边缘计算中阻塞作业调度与冗余计算资源分配，提升系统吞吐与可靠性。",
    "CacheMoE Task-Aware Expert Model Caching for Multitask Inference in Distributed Edge IoT Networks": "提出 CacheMoE，按任务感知缓存专家模型，加速分布式边缘 IoT 多任务推理。",
    "Cloud-Edge Cooperative MPC With Event-Triggered Strategy for Large-Scale Complex Systems": "面向大规模复杂系统，提出事件触发策略下的云边协同模型预测控制（MPC）。",
    "Cloud–Edge Architecture Spatiotemporal-Enhanced Infrared Target Tracking Framework in Internet of Things": "在 IoT 云边架构下，构建时空增强的红外目标跟踪框架。",
    # --- 7.31（文件名含空格/en-dash，键须与 pdf.stem 完全一致）---
    "Joint Computation Offloading and Resource Allocation for Maritime MEC With Energy Harvesting": "面向海上 MEC，研究能量收集场景下的联合计算卸载与资源分配。",
    "Joint Computation Offloading and Service Caching in Mobile Edge-Cloud Computing via Deep Reinforcement Learning": "用深度强化学习联合优化移动边云计算中的计算卸载与服务缓存。",
    "Joint Optimization of Task Offloading and Resource Allocation in Satellite-Assisted IoT Networks": "在卫星辅助 IoT 网络中，联合优化任务卸载与资源分配以提升系统性能。",
    "LsiA3CS Deep-Reinforcement-Learning-Based Cloud–Edge Collaborative Task Scheduling in Large-Scale IIoT": "提出 LsiA3CS，用深度强化学习实现大规模工业 IoT 的云边协同任务调度。",
    "Minimization of Task Completion Time in Wireless Powered Mobile Edge–Cloud Computing Networks": "在无线供电移动边云网络中，最小化任务完成时间并优化卸载决策。",
    "Optimization of Edge–Cloud Collaborative Computing Resource Management for Internet of Vehicles Based on Multiagent Deep": "基于多智能体深度方法，优化车联网边云协同计算资源管理。",
    "Optimized CNN Architectures Benchmarking in Hardware-Constrained Edge Devices in IoT Environments": "对硬件受限边缘设备上的优化 CNN 架构进行基准评测，指导 IoT 部署选型。",
    "PLAYS Minimizing DNN Inference Latency in Serverless Edge Cloud for Artificial Intelligence of Things": "提出 PLAYS，降低无服务器边云中 AIoT 场景的 DNN 推理时延。",
    "Privacy Set Privacy-Authority-Aware Compiler for Homomorphic Encryption on Edge-Cloud System": "提出 Privacy Set，面向边云同态加密的隐私权限感知编译器。",
    "Privacy-Preserving Machine Learning in Cloud–Edge–End Collaborative Environments": "研究云边端协同环境下的隐私保护机器学习方法与部署策略。",
    # --- 8.8（键须与 pdf.stem 完全一致）---
    'A Flexible and Verifiable Keyword PIR Scheme for Cloud–Edge–Terminal Collaboration in AIoT': '面向云边端 AIoT，提出可验证模糊/细粒度多关键词 PIR（B-OPRF/BF-OPRF），保护检索隐私并防云端篡改。',
    'Cloud-Edge–Terminal Collaboration-Enabled Device-Free Sensing Under Class-Imbalance Conditions': '面向云边端无源感知中的类别不平衡，用 Monitor 监测与损失校正提升 WiFi/BVP 手势识别鲁棒性。',
    'Collaborative Policy Learning for Dynamic Scheduling Tasks in Cloud–Edge–Terminal IoT Networks Using Federated Reinforce': '提出联邦强化学习协作策略学习框架，以边缘无关策略与公平任务选择加速云边端动态调度。',
    'Correlation Anomaly Detection With Multiple Primary Attributes in Collaborative Device–Edge–Cloud Network': '在端边云协同网络中，融合多主属性相关异常检测：边侧预检传播、云侧次要属性最终裁定。',
    'Many-Objective Optimization-Based Content Popularity Prediction for Cache-Assisted Cloud–Edge–End Collaborative IoT Netw': '面向缓存辅助云边端协作 IoT，用多目标优化与 KMaOEA 做内容流行度预测驱动协作缓存。',
    'MLM-WR A Swarm-Intelligence-Based Cloud–Edge–Terminal Collaboration Data Collection Scheme in the Era of AIoT': '提出 MLM-WR 群体智能云边端数据采集方案，结合 UAV-GTD 与 PSO 提升工人匹配与数据质量。',
    'Privacy-Preserving Machine Learning in Cloud–Edge–End Collaborative Environments': '研究云边端协同环境下的隐私保护机器学习方法与部署策略。',
    'QuAsyncFL Asynchronous Federated Learning With Quantization for Cloud–Edge–Terminal Collaboration Enabled AIoT': '提出 QuAsyncFL：异步联邦学习结合无偏量化，提升云边端 AIoT 联邦训练通信效率与收敛速度。',
    'Task-Oriented Multimodal Communication Based on Cloud–Edge–UAV Collaboration': '面向云边–UAV 协作的任务导向多模态通信，提升干扰/窃听下的任务成功概率。',
    'Toward Secure and Lightweight Data Transmission for Cloud–Edge–Terminal Collaboration in Artificial Intelligence of Thin': '提出 GS-SNC 轻量安全传输方案，用 Gold 序列与安全网络编码保障云边端 AIoT 数据传输安全。',

    # --- 8.16（键须与 pdf.stem 完全一致）---
    'Adaptive Model Partitioning and Pruning for Collaborative DNN Inference in Mobile Edge-Cloud Computing Networks': '面向移动边云计算，自适应模型划分与剪枝实现协作 DNN 推理加速。',
    'Adaptive Search and Collaborative Offloading Under Device-to-Device Joint Edge Computing Network': '在 D2D 联合边缘计算网络中，研究自适应搜索与协作任务卸载。',
    'CADEC A Combinatorial Auction for Dynamic Distributed DNN Inference Scheduling in Edge-Cloud Networks': '提出 CADEC：用组合拍卖为边云网络动态调度分布式 DNN 推理。',
    'CollaboRadio A Hybrid Device-Edge-Cloud Collaboration Paradigm for Fine-Grained Radio Map Construction': '提出 CollaboRadio 端边云混合协作范式，支撑细粒度无线电图构建。',
    'Federated Knowledge Distillation Using Hierarchical Reinforcement Learning in Resource-Constrained IoT Edge-Cloud Comput': '在资源受限 IoT 边云中，用分层强化学习做联邦知识蒸馏。',
    'i-CU Intelligent Cache Replacement and Content Update for Data Freshness in Cloud-Edge Networks': '提出 i-CU：云边网络中面向数据新鲜度的智能缓存替换与内容更新。',
    'Layer-Aware Cost-Effective Container Updates With Edge-Cloud Collaboration in Edge Computing': '面向边缘计算，研究层感知、成本优化的边云协作容器更新。',
    'MATE A D2D-Enhanced Multi-Bitrate Video Caching Strategy for Cloud-Edge-Device Collaborative Networks': '提出 MATE：D2D 增强的多码率视频缓存策略，服务云边端协作网络。',
    'MIDDLE A Mobility-Driven Device-Edge-Cloud Federated Learning Framework': '提出 MIDDLE：移动性驱动的端边云联邦学习框架。',
    'Symmetric Orchestration Under Service Mesh Paradigm Empowering Massive Online Applications in Edge Clouds': '在服务网格范式下研究对称编排，支撑边缘云海量在线应用。',

    # --- 8.23（键须与 pdf.stem 完全一致）---
    '01_Energy-efficient computation offloading and resource allocation in cloud-edge collabor': '提出休眠感知的云边能效卸载与资源分配（CORA），用混合动作空间多智能体 DDPG 权衡时延与能耗。',
    '02_Joint Scheduling Mechanism for Dynamic Slice Resource Allocation and Task Offloading i': '面向用户–边缘–云切片系统，用增强 MADDPG + DQN–Lagrange 做片间配额与片内卸载联合调度。',
    '03_Privacy-Preserving Task Offloading Scheme Based on Blockchain Federated Learning in En': '在端边云环境中，基于区块链联邦学习设计隐私保护任务卸载方案。',
    '04_MADRL-Based Model Partitioning, Aggregation Control, and Resource Allocation for Cloud': '面向云边端协同分裂联邦学习，用 MADRL 联合优化模型切分、聚合控制与资源分配。',
    '05_GT-MARL Graph- and Transformer-Enhanced Multi-Agent Reinforcement Learning for Cloud-E': '提出 GT-MARL：图与 Transformer 增强的多智能体强化学习云边协同调度框架。',
    '06_Game-Theoretic Bandwidth Allocation and Task Offloading in Cloud–Edge Collaboration': '用博弈论方法联合优化云边协作中的带宽分配与任务卸载决策。',
    '07_Reinforcement learning based offloading and resource allocation for multi-intelligent ': '面向绿色边云多智能车，用强化学习做实时卸载与资源分配（RT-MADDPG）。',
    '08_Joint Adaptive Aggregation and Resource Allocation for Hierarchical Federated Learning': '面向边云协作分层联邦学习，联合自适应聚合频率与资源分配（AHFLP）。',
    '09_Task Offloading and Resource Scheduling in Mobile Edge-Cloud Computing Based on Edge C': '在移动边云计算中，结合边缘竞争与任务预测进行卸载与资源调度。',
    '10_Adversarial Bandit Learning Assisted Online Optimization for Digital Twin Placement an': '用对抗老虎机在线优化端边云数字孪生放置与更新（ARBOK）。',

}


def title_from_stem(stem: str) -> str:
    """将文件名转为可读英文标题。"""
    return stem.replace("_", " ").replace("-", " ").strip()


def fallback_summary(stem: str) -> str:
    """无预置摘要时，根据标题生成简短说明。"""
    title = title_from_stem(stem)
    return f"研究 {title.lower()} 相关问题的学术论文。"


# 与更早批次重复的 PDF：保留首见，跳过后续副本（仍留在磁盘）
# 用 (日期文件夹, 文件名 stem 前缀) 匹配，避免 en-dash 字符不一致
SKIP_STEM_PREFIXES = {
    ("8.8", "Privacy-Preserving Machine Learning in Cloud"),
}


def scan_literature() -> list[dict]:
    items = []
    for folder in DATE_FOLDERS:
        folder_path = ROOT / folder
        if not folder_path.is_dir():
            continue
        for pdf in sorted(folder_path.glob("*.pdf")):
            stem = pdf.stem
            if any(folder == f and stem.startswith(pref) for f, pref in SKIP_STEM_PREFIXES):
                continue
            rel = str(pdf.relative_to(ROOT)).replace("\\", "/")
            items.append(
                {
                    "date_folder": folder,
                    "filename": pdf.name,
                    "title": title_from_stem(stem),
                    "summary": SUMMARIES.get(stem, fallback_summary(stem)),
                    "path": str(pdf.resolve()),
                    "relative_path": rel,
                }
            )
    return items


def build_ppt(items: list[dict], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "文献速览 · 云边协同研究"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x47, 0x8A)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"共 {len(items)} 篇 · 按阅读日期分组 · 点击标题可打开 PDF"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    # 总目录
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(6.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "文献总目录"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x47, 0x8A)

    current_folder = None
    for item in items:
        if item["date_folder"] != current_folder:
            current_folder = item["date_folder"]
            hp = tf.add_paragraph()
            hp.text = f"\n📅 {current_folder}"
            hp.font.size = Pt(18)
            hp.font.bold = True
            hp.font.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
        lp = tf.add_paragraph()
        lp.text = f"  • {item['title'][:80]}{'…' if len(item['title']) > 80 else ''}"
        lp.font.size = Pt(14)
        lp.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        lp.level = 1

    # 按日期分组，每页 2 篇
    grouped: dict[str, list] = {}
    for item in items:
        grouped.setdefault(item["date_folder"], []).append(item)

    for folder, papers in grouped.items():
        for i in range(0, len(papers), 2):
            batch = papers[i : i + 2]
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            header = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6))
            hp = header.text_frame.paragraphs[0]
            hp.text = f"📅 {folder}  ·  文献摘要"
            hp.font.size = Pt(24)
            hp.font.bold = True
            hp.font.color.rgb = RGBColor(0x1A, 0x47, 0x8A)

            for idx, paper in enumerate(batch):
                top = 1.1 + idx * 3.0
                card = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(2.7))
                tf = card.text_frame
                tf.word_wrap = True

                tp = tf.paragraphs[0]
                tp.text = paper["title"]
                tp.font.size = Pt(16)
                tp.font.bold = True
                tp.font.color.rgb = RGBColor(0x15, 0x65, 0xC0)

                sp = tf.add_paragraph()
                sp.text = paper["summary"]
                sp.font.size = Pt(14)
                sp.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
                sp.space_before = Pt(8)

                fp = tf.add_paragraph()
                fp.text = f"📄 {paper['filename']}"
                fp.font.size = Pt(11)
                fp.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
                fp.space_before = Pt(6)

                # 超链接到 PDF
                slide.shapes.add_shape(1, Inches(0.5), Inches(top), Inches(12.3), Inches(2.7)).click_action.hyperlink.address = paper["path"]

    prs.save(output_path)


def build_html(items: list[dict], output_path: Path) -> None:
    data_json = json.dumps(items, ensure_ascii=False)
    zones_path = OUTPUT_DIR / "journal_zones.json"
    zones_json = json.dumps(json.loads(zones_path.read_text(encoding="utf-8")), ensure_ascii=False)
    html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>文献检索 · EdgeNexus 边云合智</title>
<link rel="icon" href="assets/favicon.svg" type="image/svg+xml">
<link rel="apple-touch-icon" href="assets/favicon.png">
<style>
  :root {{
    --bg: #f0f4f8; --card: #fff; --primary: #1a478a; --accent: #2e7d32;
    --text: #333; --muted: #888; --border: #dde3ea; --highlight: #fff3cd;
  }}
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{ font-family: "Segoe UI", "Microsoft YaHei", sans-serif; background: var(--bg); color: var(--text); min-height: 100vh; }}
  .site-nav {{
    position: sticky; top: 0; z-index: 100;
    display: flex; align-items: center; justify-content: space-between;
    padding: 0 28px; height: 52px;
    background: #0f2f5c; border-bottom: 1px solid rgba(255,255,255,.08);
  }}
  .site-nav .brand {{
    display: inline-flex; align-items: center; gap: 10px;
    color: #fff; font-weight: 700; font-size: .98rem; text-decoration: none;
  }}
  .site-nav .brand img {{ width: 28px; height: 28px; border-radius: 7px; }}
  .site-nav .links {{ display: flex; gap: 6px; flex-wrap: wrap; }}
  .site-nav .links a {{
    color: rgba(255,255,255,.78); text-decoration: none; font-size: .9rem; font-weight: 600;
    padding: 8px 14px; border-radius: 8px; transition: background .15s, color .15s;
  }}
  .site-nav .links a:hover {{ background: rgba(255,255,255,.1); color: #fff; }}
  .site-nav .links a.active {{ background: rgba(255,255,255,.16); color: #fff; }}
  header {{ background: linear-gradient(135deg, #1a478a, #2563b0); color: #fff; padding: 28px 32px; }}
  header h1 {{ font-size: 1.6rem; font-weight: 700; }}
  header p {{ opacity: .85; margin-top: 6px; font-size: .95rem; }}
  .search-bar {{ padding: 20px 32px; background: var(--card); border-bottom: 1px solid var(--border); position: sticky; top: 52px; z-index: 10; box-shadow: 0 2px 8px rgba(0,0,0,.06); }}
  .search-bar input {{ width: 100%; padding: 14px 18px; font-size: 1rem; border: 2px solid var(--border); border-radius: 10px; outline: none; transition: border .2s; }}
  .search-bar input:focus {{ border-color: var(--primary); }}
  .page-layout {{ display: flex; align-items: flex-start; max-width: 1180px; margin: 0 auto; padding: 0 20px; gap: 24px; }}
  .date-nav {{ position: sticky; top: 140px; width: 108px; flex-shrink: 0; max-height: calc(100vh - 148px); overflow-y: auto; padding: 8px 6px 16px 0; scrollbar-width: thin; scrollbar-color: #c5cfe0 transparent; }}
  .date-nav::-webkit-scrollbar {{ width: 5px; }}
  .date-nav::-webkit-scrollbar-thumb {{ background: #c5cfe0; border-radius: 999px; }}
  .date-nav-title {{ font-size: .78rem; font-weight: 700; color: var(--primary); margin-bottom: 10px; letter-spacing: .02em; position: sticky; top: 0; background: var(--bg); padding-bottom: 6px; z-index: 1; }}
  .date-nav-list {{ list-style: none; display: flex; flex-direction: column; gap: 4px; }}
  .date-nav-item {{ display: block; padding: 7px 10px; border-radius: 8px; font-size: .82rem; font-weight: 600; color: #555; text-decoration: none; border: 1px solid transparent; transition: background .15s, color .15s, border-color .15s; cursor: pointer; }}
  .date-nav-item:hover {{ background: #e8eef5; color: var(--primary); }}
  .date-nav-item.active {{ background: #e3f2fd; color: var(--primary); border-color: #bbdefb; }}
  .date-nav-item.disabled {{ opacity: .35; pointer-events: none; }}
  .date-nav-count {{ display: block; font-size: .68rem; font-weight: 500; color: var(--muted); margin-top: 1px; }}
  .page-main {{ flex: 1; min-width: 0; }}
  .stats {{ padding: 8px 0 0; font-size: .85rem; color: var(--muted); }}
  .container {{ padding: 16px 0 40px; }}
  .date-group {{ margin-top: 24px; scroll-margin-top: 96px; }}
  .date-label {{ font-size: .9rem; font-weight: 700; color: var(--accent); margin-bottom: 10px; padding-left: 4px; border-left: 4px solid var(--accent); padding-left: 10px; }}
  .card {{ background: var(--card); border: 1px solid var(--border); border-radius: 12px; padding: 16px 18px; margin-bottom: 10px; transition: box-shadow .2s, transform .15s; cursor: pointer; }}
  .card:hover {{ box-shadow: 0 4px 16px rgba(26,71,138,.12); transform: translateY(-1px); }}
  .card.hidden {{ display: none; }}
  .card-title {{ font-size: .95rem; font-weight: 600; color: var(--primary); line-height: 1.4; }}
  .card-summary {{ font-size: .88rem; color: #555; margin-top: 6px; line-height: 1.55; }}
  .card-meta {{ font-size: .78rem; color: var(--muted); margin-top: 8px; }}
  .card mark {{ background: var(--highlight); border-radius: 2px; padding: 0 2px; }}
  .no-result {{ text-align: center; padding: 60px 20px; color: var(--muted); display: none; }}
  .no-result.show {{ display: block; }}
  .open-btn {{ display: inline-block; margin-top: 8px; padding: 5px 12px; background: var(--primary); color: #fff; border-radius: 6px; font-size: .78rem; text-decoration: none; }}
  .open-btn:hover {{ background: #2563b0; }}
  .open-btn.translate {{ background: var(--accent); margin-left: 6px; }}
  .open-btn.translate:hover {{ background: #388e3c; }}
  .lookup-panel {{ margin: 12px 0 0; padding: 14px 16px; background: #fff; border: 1px solid var(--border); border-radius: 10px; display: none; }}
  .lookup-panel.show {{ display: block; }}
  .lookup-panel h3 {{ font-size: .92rem; color: var(--primary); margin-bottom: 8px; }}
  .lookup-row {{ font-size: .86rem; line-height: 1.7; color: #444; }}
  .badges {{ display: flex; flex-wrap: wrap; gap: 6px; margin-top: 8px; }}
  .badge {{ display: inline-block; padding: 2px 8px; border-radius: 999px; font-size: .72rem; font-weight: 600; }}
  .badge-local-yes {{ background: #e8f5e9; color: #2e7d32; }}
  .badge-local-no {{ background: #fff3e0; color: #e65100; }}
  .badge-zone1 {{ background: #e3f2fd; color: #1565c0; }}
  .badge-not-zone1 {{ background: #fce4ec; color: #c62828; }}
  .badge-unknown {{ background: #f5f5f5; color: #757575; }}
  .badge-cecc-yes {{ background: #e8eaf6; color: #283593; }}
  .badge-cecc-no {{ background: #efebe9; color: #5d4037; }}
  .card-journal {{ font-size: .78rem; color: #666; margin-top: 6px; }}
</style>
</head>
<body>
<nav class="site-nav" aria-label="站点导航">
  <a class="brand" href="index.html">
    <img src="assets/favicon.svg" alt="">
    <span>EdgeNexus · 边云合智</span>
  </a>
  <div class="links">
    <a href="index.html" class="active">文献检索</a>
  </div>
</nav>
<header>
  <h1>📚 文献检索 · 云边协同研究</h1>
  <p>输入论文文件名或英文标题 · 自动判断本地是否已有、是否中科院一区、是否云边协同计算方向 · 点击卡片打开 PDF</p>
</header>
<div class="search-bar">
  <input id="searchInput" type="text" placeholder="例如：Flocky.pdf / task offloading / Cloud-Edge_Collaborative …" autofocus>
</div>
<div class="page-layout">
  <nav class="date-nav" id="dateNav" aria-label="阅读日期目录"></nav>
  <div class="page-main">
    <div class="lookup-panel" id="lookupPanel"></div>
    <p class="stats" id="stats"></p>
    <div class="container" id="results"></div>
    <div class="no-result" id="noResult"><p>未找到匹配的文献，请换个关键词试试</p></div>
  </div>
</div>

<script>
const PAPERS = {data_json};
const JOURNAL_ZONES = {zones_json};

function normalize(s) {{
  return s.toLowerCase()
    .replace(/[_\\-–—]/g, ' ')
    .replace(/[^\\w\\s\\u4e00-\\u9fff]/g, ' ')
    .replace(/\\s+/g, ' ')
    .trim();
}}

function paperHaystack(p) {{
  return normalize([
    p.title,
    p.crossref_title,
    p.summary,
    p.filename,
    p.journal,
    p.date_folder,
  ].filter(Boolean).join(' '));
}}

function paperMatchesQuery(p, q) {{
  if (!q) return true;
  const qNorm = normalize(q);
  const hay = paperHaystack(p);
  if (!qNorm) return true;
  if (hay.includes(qNorm)) return true;
  const words = qNorm.split(' ').filter(w => w.length > 1);
  if (!words.length) return hay.includes(qNorm);
  return words.every(w => hay.includes(w));
}}

function ceccScore(norm) {{
  if (!norm) return {{ score: 0, detail: '无有效文本' }};

  const core = [
    /cloud[\s\-–—]?edge/i, /edge[\s\-–—]?cloud/i, /cloudedge/i, /edgecloud/i,
    /cloud[\s\-–—]?edge[\s\-–—]?(end|terminal)/i, /end[\s\-–—]?edge[\s\-–—]?cloud/i,
    /edge[\s\-–—]?hub[\s\-–—]?cloud/i, /edge[\s\-–—]?cloud[\s\-–—]?continuum/i,
    /cloud[\s\-–—]?edge[\s\-–—]?(collabor|cooper)/i, /edge[\s\-–—]?cloud[\s\-–—]?(collabor|cooper)/i,
    /edge[\s\-–—]?collabor/i, /mobile[\s\-–—]?cloud[\s\-–—]?edge/i, /serverless[\s\-–—]?edge/i,
    /云边协同/i, /云边协作/i, /云边计算/i, /云边/i, /边云/i,
  ];
  if (core.some(re => re.test(norm))) return {{ score: 5, detail: '匹配云边协同核心主题' }};

  const zh = [/边缘协作/, /边缘计算/, /边计算/, /边缘编排/, /无服务器边/, /云边/, /边云/, /卸载/, /调度/, /端边云/, /边云协同/, /云边协同/];
  if (zh.some(re => re.test(norm))) return {{ score: 4, detail: '匹配云边协同中文主题' }};

  let score = 0;
  const reasons = [];
  const edge = [
    /mobile[\s\-–—]?edge[\s\-–—]?comput/i, /multi[\s\-–—]?access[\s\-–—]?edge/i, /\bmec\b/i,
    /edge[\s\-–—]?comput/i, /vehicular[\s\-–—]?edge/i, /uav[\s\-–—]?(based[\s\-–—]?)?edge/i,
    /edge[\s\-–—]?server/i, /edge[\s\-–—]?node/i, /collaborative[\s\-–—]?edge[\s\-–—]?comput/i,
    /serverless[\s\-–—]?edge/i, /federated[\s\-–—]?edge[\s\-–—]?learn/i,
    /edge[\s\-–—]?orchestrat/i, /edge[\s\-–—]?shard/i, /edge[\s\-–—]?collabor/i,
    /dynamic[\s\-–—]?edge[\s\-–—]?network/i, /\bon[\s\-–—]?edge\b/i, /tasks[\s\-–—]?on[\s\-–—]?edge/i,
  ];
  const task = [
    /task[\s\-–—]?offload/i, /computation[\s\-–—]?offload/i, /compute[\s\-–—]?offload/i,
    /computing[\s\-–—]?offload/i, /resource[\s\-–—]?alloc/i,
    /service[\s\-–—]?(cach|deploy|migr|orchestr|placement|select|routing)/i,
    /workflow[\s\-–—]?(schedul|container)/i, /task[\s\-–—]?(placement|schedul|alloc)/i,
    /load[\s\-–—]?balanc/i, /container/i, /serverless/i, /\bsfc\b/i, /kubernetes|\bk8s\b/i,
  ];
  const method = [
    /federated[\s\-–—]?learn/i, /federated[\s\-–—]?reinforcement/i, /deep[\s\-–—]?reinforcement/i,
    /\bdrl\b/i, /\bmarl\b/i, /\bmaddpg\b/i, /lyapunov/i, /blockchain/i, /digital[\s\-–—]?twin/i,
    /auction/i, /game[\s\-–—]?theoret/i, /stackelberg/i, /differential[\s\-–—]?privacy/i,
    /over[\s\-–—]?the[\s\-–—]?air/i, /beamform/i, /client[\s\-–—]?sampl/i,
  ];

  const hasEdge = edge.some(re => re.test(norm));
  const hasTask = task.some(re => re.test(norm));
  const hasCloud = /\bcloud\b/i.test(norm);
  const hasMethod = method.some(re => re.test(norm));
  const hasCollab = /collabor|cooper/i.test(norm);
  const hasCloudAssist = /cloud[\s\-–—]?assist/i.test(norm);
  const hasInfer = /infer|inference|llm|dnn|deep[\s\-–—]?learn/i.test(norm);
  const hasIot = /\biot\b|internet[\s\-–—]?of[\s\-–—]?things|iiot|aiot/i.test(norm);

  if (hasEdge) {{ score += 3; reasons.push('边缘计算场景'); }}
  if (hasTask) {{ score += 2; reasons.push('卸载/调度/编排'); }}
  if (hasCloud && (hasEdge || hasCollab || hasTask)) {{ score += 3; reasons.push('云-边联合优化'); }}
  if (hasCloudAssist && hasEdge) {{ score += 3; reasons.push('云辅助边端计算'); }}
  if (hasCollab && (hasEdge || hasInfer || hasCloud)) {{ score += 2; reasons.push('多层协作计算/推理'); }}
  if (hasInfer && hasEdge) {{ score += 2; reasons.push('边缘 AI 推理'); }}
  if (/v2x|vehicular/i.test(norm) && (hasTask || hasEdge)) {{ score += 2; reasons.push('车联网边缘优化'); }}
  if (/federated[\s\-–—]?learn/i.test(norm) && /edge|iot|mec|iiot|aiot|aggreg|device|client|communic|privacy|offload|resource|over[\s\-–—]?the[\s\-–—]?air|beamform|sampling/i.test(norm)) {{
    score += 2; reasons.push('云边协同方法学(联邦学习/边缘通信)');
  }} else if (hasMethod && (hasEdge || hasIot || hasTask || hasCloud)) {{
    score += 2; reasons.push('云边协同方法学(优化/DRL/机制等)');
  }}
  if (/satellite/i.test(norm) && hasEdge) {{ score += 2; reasons.push('空天地边缘协同'); }}
  if (/energy|latency|delay|aoi|reliability|trust|privacy|cost|budget/i.test(norm) && (hasEdge || hasTask || hasCloud)) {{
    score += 1; reasons.push('CECC 典型优化目标');
  }}

  return {{ score, detail: reasons.slice(0, 3).join(' · ') || '未识别云边协同计算相关主题' }};
}}

function classifyCloudEdgeResearch(text, localCorpus = false) {{
  const norm = normalize(String(text || '').replace(/_/g, ' ').replace(/-/g, ' '));
  const {{ score, detail }} = ceccScore(norm);
  if (localCorpus) {{
    return {{
      is_cecc: true,
      cecc_label: '是',
      cecc_detail: '本地精读库 · ' + (score > 0 ? detail : '云边协同精读文献'),
    }};
  }}
  if (score >= 2) return {{ is_cecc: true, cecc_label: '是', cecc_detail: detail }};
  return {{ is_cecc: false, cecc_label: '否', cecc_detail: detail }};
}}

function ceccFromPaper(p) {{
  const text = [p.title, p.crossref_title, p.summary, p.filename].filter(Boolean).join(' ');
  if (p.local) return classifyCloudEdgeResearch(text, true);
  if (p.is_cecc != null) {{
    return {{ is_cecc: p.is_cecc, cecc_label: p.cecc_label || (p.is_cecc ? '是' : '否'), cecc_detail: p.cecc_detail || '' }};
  }}
  return classifyCloudEdgeResearch(text, false);
}}

function ceccBadge(isCecc, label) {{
  if (isCecc === true) return `<span class="badge badge-cecc-yes">云边协同：${{label || '是'}}</span>`;
  if (isCecc === false) return `<span class="badge badge-cecc-no">云边协同：${{label || '否'}}</span>`;
  return `<span class="badge badge-unknown">云边协同：${{label || '未知'}}</span>`;
}}

function matchJournal(journalName) {{
  const journals = JOURNAL_ZONES.journals || {{}};
  const aliases = JOURNAL_ZONES.aliases || {{}};
  if (!journalName) return {{ matched: null, info: null }};

  let raw = journalName.trim();
  if (aliases[raw]) raw = aliases[raw];
  if (journals[raw]) return {{ matched: raw, info: journals[raw] }};

  const norm = normalize(raw);
  for (const [key, info] of Object.entries(journals)) {{
    if (normalize(key) === norm) return {{ matched: key, info }};
  }}
  for (const [alias, target] of Object.entries(aliases)) {{
    if (normalize(alias) === norm) return {{ matched: target, info: journals[target] }};
  }}
  for (const [key, info] of Object.entries(journals)) {{
    const kn = normalize(key);
    if (kn.includes(norm) || norm.includes(kn)) return {{ matched: key, info }};
  }}
  return {{ matched: raw, info: null }};
}}

function zoneFromJournal(journalName) {{
  const {{ matched, info }} = matchJournal(journalName);
  if (!info) {{
    return {{
      journal: journalName || '',
      cas_zone: null,
      jcr_quartile: null,
      is_zone1: null,
      zone_label: '分区未知',
      zone_detail: '期刊未收录在本地分区库，可手动核对中科院分区',
    }};
  }}
  const cas = info.cas_zone;
  const jcr = info.jcr_quartile;
  const isZone1 = cas === 1;
  const label = isZone1 ? '中科院一区' : (cas != null ? `中科院${{cas}}区` : '分区未知');
  const parts = [];
  if (cas != null) parts.push(`中科院 ${{cas}} 区`);
  if (jcr) parts.push(`JCR ${{jcr}}`);
  return {{
    journal: journalName,
    matched_journal: matched,
    cas_zone: cas,
    jcr_quartile: jcr,
    is_zone1: isZone1,
    zone_label: label,
    zone_detail: parts.join(' · '),
  }};
}}

function titleScore(query, crTitle) {{
  const a = normalize(query);
  const b = normalize(crTitle || '');
  if (!a || !b) return 0;
  if (a === b) return 1;
  if (a.includes(b) || b.includes(a)) {{
    return 0.7 + 0.15 * (Math.min(a.length, b.length) / Math.max(a.length, b.length));
  }}
  const aw = new Set(a.split(/\\s+/));
  const bw = new Set(b.split(/\\s+/));
  let hit = 0;
  aw.forEach(w => {{ if (bw.has(w)) hit++; }});
  const wordScore = hit / aw.size;
  if (aw.size <= 2 && wordScore < 1) return wordScore * 0.5;
  return wordScore;
}}

async function fetchJson(url) {{
  const resp = await fetch(url);
  if (!resp.ok) throw new Error('HTTP ' + resp.status);
  return resp.json();
}}

async function searchOpenAlex(title) {{
  const url = 'https://api.openalex.org/works?search=' + encodeURIComponent(title.slice(0, 120)) + '&per_page=5';
  const data = await fetchJson(url);
  const items = data.results || [];
  if (!items.length) return null;

  let best = items[0];
  let bestScore = titleScore(title, best.title);
  for (const item of items) {{
    const s = titleScore(title, item.title);
    if (s > bestScore) {{ best = item; bestScore = s; }}
  }}
  if (bestScore < 0.45) return null;

  const journal = best.primary_location?.source?.display_name || '';
  return {{
    crossref_title: best.title,
    journal,
    doi: best.doi || '',
    match_score: Math.round(bestScore * 100) / 100,
  }};
}}

async function searchCrossref(title) {{
  const url = 'https://api.crossref.org/works?query.title=' + encodeURIComponent(title.slice(0, 120)) + '&rows=5';
  const data = await fetchJson(url);
  const allowed = new Set(['journal-article', 'proceedings-article']);
  let items = (data.message?.items || []).filter(it => allowed.has(it.type));
  if (!items.length) items = data.message?.items || [];
  if (!items.length) return null;

  let best = items[0];
  let bestScore = titleScore(title, (best.title || [''])[0]);
  for (const item of items) {{
    const s = titleScore(title, (item.title || [''])[0]);
    if (s > bestScore) {{ best = item; bestScore = s; }}
  }}
  if (bestScore < 0.45) return null;

  const journal = (best['container-title'] || [])[0] || '';
  return {{
    crossref_title: (best.title || [''])[0],
    journal,
    doi: best.DOI || '',
    match_score: Math.round(bestScore * 100) / 100,
  }};
}}

async function fetchZoneLookup(title) {{
  const empty = {{
    query: title,
    found: false,
    local: false,
    crossref_title: '',
    journal: '',
    is_zone1: null,
    zone_label: '未找到论文',
    zone_detail: '未匹配到该标题，请检查拼写或换用英文全称',
  }};
  try {{
    let hit = null;
    try {{
      hit = await searchOpenAlex(title);
    }} catch (_) {{}}
    if (!hit) {{
      try {{
        hit = await searchCrossref(title);
      }} catch (_) {{}}
    }}
    if (!hit) return empty;

    const zone = zoneFromJournal(hit.journal);
    const cecc = classifyCloudEdgeResearch(title + ' ' + hit.crossref_title);
    return {{
      query: title,
      found: true,
      local: false,
      crossref_title: hit.crossref_title,
      journal: hit.journal,
      doi: hit.doi,
      match_score: hit.match_score,
      ...zone,
      ...cecc,
    }};
  }} catch (e) {{
    return {{
      ...empty,
      zone_label: '无法查询分区',
      zone_detail: '请确认已联网；若仍失败可稍后再试',
    }};
  }}
}}

function highlight(text, query) {{
  if (!query) return text;
  const re = new RegExp('(' + query.replace(/[.*+?^${{}}()|[\\]\\\\]/g, '\\\\$&') + ')', 'gi');
  return text.replace(re, '<mark>$1</mark>');
}}

function openPdf(path) {{
  window.open('file:///' + path.replace(/\\\\/g, '/'), '_blank');
}}

function zoneBadge(isZone1, label) {{
  if (isZone1 === true) return `<span class="badge badge-zone1">✓ ${{label || '中科院一区'}}</span>`;
  if (isZone1 === false) return `<span class="badge badge-not-zone1">✗ ${{label || '非一区'}}</span>`;
  return `<span class="badge badge-unknown">${{label || '分区未知'}}</span>`;
}}

function localBadge(hasLocal) {{
  return hasLocal
    ? '<span class="badge badge-local-yes">✓ 本地已有</span>'
  : '<span class="badge badge-local-no">✗ 本地未下载</span>';
}}

function renderLookupPanel(q, localMatches, external) {{
  const panel = document.getElementById('lookupPanel');
  if (!q) {{
    panel.classList.remove('show');
    panel.innerHTML = '';
    return;
  }}

  const hasLocal = localMatches.length > 0;
  let zoneHtml = '';
  let ceccHtml = '';
  let journalHtml = '';
  let titleHtml = '';
  let ceccDetailHtml = '';

  if (hasLocal) {{
    const p = localMatches[0];
    const cecc = ceccFromPaper(p);
    zoneHtml = zoneBadge(p.is_zone1, p.zone_label);
    ceccHtml = ceccBadge(cecc.is_cecc, cecc.cecc_label);
    journalHtml = p.journal ? `<div class="lookup-row">📰 期刊：${{p.journal}}</div>` : '';
    titleHtml = `<div class="lookup-row">📄 匹配：${{p.title}}</div>`;
    if (cecc.cecc_detail) ceccDetailHtml = `<div class="lookup-row">🔬 方向判定：${{cecc.cecc_detail}}</div>`;
  }} else if (external) {{
    if (external.found) {{
      zoneHtml = zoneBadge(external.is_zone1, external.zone_label);
      ceccHtml = ceccBadge(external.is_cecc, external.cecc_label);
      journalHtml = external.journal ? `<div class="lookup-row">📰 期刊：${{external.journal}}</div>` : '';
      titleHtml = external.crossref_title ? `<div class="lookup-row">🔍 Crossref：${{external.crossref_title}}</div>` : '';
      if (external.cecc_detail) ceccDetailHtml = `<div class="lookup-row">🔬 方向判定：${{external.cecc_detail}}</div>`;
    }} else {{
      zoneHtml = `<span class="badge badge-unknown">${{external.zone_label || '查询中…'}}</span>`;
      const cecc = classifyCloudEdgeResearch(q);
      ceccHtml = ceccBadge(cecc.is_cecc, cecc.cecc_label);
      if (external.zone_detail) journalHtml = `<div class="lookup-row">${{external.zone_detail}}</div>`;
      if (cecc.cecc_detail) ceccDetailHtml = `<div class="lookup-row">🔬 方向判定：${{cecc.cecc_detail}}</div>`;
    }}
  }} else {{
      zoneHtml = q.length < 8
        ? '<span class="badge badge-unknown">请输入更多关键词以查询分区</span>'
        : '<span class="badge badge-unknown">分区查询中…</span>';
      const cecc = classifyCloudEdgeResearch(q);
      ceccHtml = ceccBadge(cecc.is_cecc, cecc.cecc_label);
      if (cecc.cecc_detail) ceccDetailHtml = `<div class="lookup-row">🔬 方向判定：${{cecc.cecc_detail}}</div>`;
  }}

  panel.innerHTML = `
    <h3>论文速查 · 「${{q}}」</h3>
    <div class="badges">${{localBadge(hasLocal)}} ${{zoneHtml}} ${{ceccHtml}}</div>
    ${{titleHtml}}${{journalHtml}}${{ceccDetailHtml}}
  `;
  panel.classList.add('show');
}}

let lookupTimer = null;

function scheduleLookup(q, localMatches) {{
  clearTimeout(lookupTimer);
  if (!q) {{
    renderLookupPanel('', [], null);
    return;
  }}
  if (localMatches.length > 0) {{
    renderLookupPanel(q, localMatches, null);
    return;
  }}
  renderLookupPanel(q, [], null);
  if (q.length < 8) return;
  lookupTimer = setTimeout(async () => {{
    const external = await fetchZoneLookup(q);
    if (normalize(document.getElementById('searchInput').value.trim()) === normalize(q)) {{
      renderLookupPanel(q, [], external);
    }}
  }}, 450);
}}

function dateAnchorId(dateFolder) {{
  return 'date-' + String(dateFolder).replace(/\\./g, '-');
}}

function getAllDates() {{
  const seen = new Set();
  const dates = [];
  PAPERS.forEach(p => {{
    if (!seen.has(p.date_folder)) {{
      seen.add(p.date_folder);
      dates.push(p.date_folder);
    }}
  }});
  return dates;
}}

function countByDate(papers) {{
  const map = {{}};
  papers.forEach(p => {{ map[p.date_folder] = (map[p.date_folder] || 0) + 1; }});
  return map;
}}

function renderDateNav(filteredPapers) {{
  const nav = document.getElementById('dateNav');
  const allDates = getAllDates();
  const activeSet = new Set((filteredPapers || PAPERS).map(p => p.date_folder));
  const counts = countByDate(filteredPapers || PAPERS);

  nav.innerHTML = `
    <div class="date-nav-title">📅 日期目录</div>
    <div class="date-nav-list">
      ${{allDates.map(d => {{
        const enabled = activeSet.has(d);
        const cls = ['date-nav-item', enabled ? '' : 'disabled'].filter(Boolean).join(' ');
        return `<a class="${{cls}}" href="#${{dateAnchorId(d)}}" data-date="${{d}}">${{d}}<span class="date-nav-count">${{counts[d] || 0}} 篇</span></a>`;
      }}).join('')}}
    </div>
  `;

  nav.querySelectorAll('.date-nav-item:not(.disabled)').forEach(link => {{
    link.addEventListener('click', (e) => {{
      e.preventDefault();
      const target = document.getElementById(dateAnchorId(link.dataset.date));
      if (target) {{
        target.scrollIntoView({{ behavior: 'smooth', block: 'start' }});
        nav.querySelectorAll('.date-nav-item').forEach(el => el.classList.remove('active'));
        link.classList.add('active');
      }}
    }});
  }});
}}

function updateActiveDateNav() {{
  const groups = [...document.querySelectorAll('.date-group')];
  if (!groups.length) return;
  const marker = window.scrollY + 110;
  let current = groups[0];
  groups.forEach(g => {{
    if (g.offsetTop <= marker) current = g;
  }});
  const date = current.dataset.date;
  document.querySelectorAll('.date-nav-item').forEach(el => {{
    el.classList.toggle('active', el.dataset.date === date);
  }});
}}

let scrollSpyTimer = null;
window.addEventListener('scroll', () => {{
  clearTimeout(scrollSpyTimer);
  scrollSpyTimer = setTimeout(updateActiveDateNav, 60);
}});

function render(query) {{
  const q = query.trim();
  const container = document.getElementById('results');
  const noResult = document.getElementById('noResult');
  container.innerHTML = '';

  const filtered = PAPERS.filter(p => paperMatchesQuery(p, q));

  document.getElementById('stats').textContent =
    q ? `找到 ${{filtered.length}} / ${{PAPERS.length}} 篇文献` : `共 ${{PAPERS.length}} 篇文献 · 按阅读日期分组`;

  scheduleLookup(q, filtered);
  renderDateNav(filtered);

  if (filtered.length === 0) {{
    noResult.classList.add('show');
    return;
  }}
  noResult.classList.remove('show');

  let currentDate = '';
  filtered.forEach(p => {{
    if (p.date_folder !== currentDate) {{
      currentDate = p.date_folder;
      const g = document.createElement('div');
      g.className = 'date-group';
      g.id = dateAnchorId(p.date_folder);
      g.dataset.date = p.date_folder;
      g.innerHTML = `<div class="date-label">📅 ${{p.date_folder}}</div>`;
      container.appendChild(g);
    }}
    const card = document.createElement('div');
    card.className = 'card';
    const cecc = ceccFromPaper(p);
    const journalLine = p.journal ? `<div class="card-journal">📰 ${{p.journal}} · ${{p.zone_detail || p.zone_label || ''}}</div>` : '';
    card.innerHTML = `
      <div class="card-title">${{highlight(p.title, q)}}</div>
      <div class="badges">${{localBadge(true)}} ${{zoneBadge(p.is_zone1, p.zone_label)}} ${{ceccBadge(cecc.is_cecc, cecc.cecc_label)}}</div>
      <div class="card-summary">${{highlight(p.summary, q)}}</div>
      ${{journalLine}}
      <div class="card-meta">📄 ${{p.filename}}</div>
      <a class="open-btn" href="file:///${{p.path.replace(/\\\\/g, '/')}}" target="_blank">打开 PDF ↗</a>
    `;
    card.onclick = (e) => {{
      if (e.target.tagName !== 'A') openPdf(p.path);
    }};
    container.lastElementChild.appendChild(card);
  }});
  updateActiveDateNav();
}}

const input = document.getElementById('searchInput');
input.addEventListener('input', () => render(input.value));
input.addEventListener('keydown', (e) => {{
  if (e.key === 'Enter') {{
    const filtered = PAPERS.filter(p => paperMatchesQuery(p, input.value.trim()));
    if (filtered.length === 1) openPdf(filtered[0].path);
  }}
}});
render('');
</script>
</body>
</html>"""
    output_path.write_text(html, encoding="utf-8")


def ensure_cecc_fields(item: dict) -> dict:
    merged = dict(item)
    merged.update(
        classify_cloud_edge_research(
            " ".join(
                filter(
                    None,
                    [item.get("title"), item.get("filename"), item.get("summary"), item.get("crossref_title")],
                )
            ),
            local_corpus=True,
        )
    )
    return merged


def enrich_items(items: list[dict], existing: list[dict] | None = None) -> list[dict]:
    by_path = {item["path"]: item for item in (existing or []) if item.get("zone_label")}
    enriched = []
    for i, item in enumerate(items, 1):
        cached = by_path.get(item["path"])
        if cached and cached.get("zone_label") and cached.get("zone_label") != "未查询":
            merged = dict(item)
            for key in ("journal", "crossref_title", "cas_zone", "jcr_quartile", "is_zone1", "zone_label", "zone_detail", "local", "is_cecc", "cecc_label", "cecc_detail"):
                if key in cached:
                    merged[key] = cached[key]
            enriched.append(merged)
            continue
        print(f"  查询分区 [{i}/{len(items)}] {item['title'][:50]}…")
        enriched.append(enrich_paper_item(item))
    return enriched


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    items = scan_literature()
    if not items:
        print("未找到 PDF 文献。")
        return

    catalog_path = OUTPUT_DIR / "catalog.json"
    existing = []
    if catalog_path.exists():
        existing = json.loads(catalog_path.read_text(encoding="utf-8"))

    new_count = len([i for i in items if i["path"] not in {e["path"] for e in existing}])
    if new_count:
        print(f"发现 {new_count} 篇新文献，正在查询期刊分区…")
    items = enrich_items(items, existing)
    items = [ensure_cecc_fields(i) for i in items]
    catalog_path.write_text(json.dumps(items, ensure_ascii=False, indent=2), encoding="utf-8")

    html_path = OUTPUT_DIR / "index.html"
    build_html(items, html_path)

    print(f"共扫描 {len(items)} 篇文献")
    print(f"目录数据: {catalog_path}")
    print(f"检索页: {html_path}")


if __name__ == "__main__":
    main()
